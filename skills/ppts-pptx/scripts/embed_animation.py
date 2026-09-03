#!/usr/bin/env python3
"""动画素材嵌入器：把 html_capture.py 的产物嵌入 PPTX 指定页（Phase 1）。

语义（混合交付 Phase 1）：
- GIF 优先：--gif 给出时嵌入动图，放映模式自动循环播放；
- 快照兜底：无 --gif 时嵌入静态快照（--snapshot）；
- 超链接常挂：给出 --html 时，图片单击跳转随附的交互 HTML 原件
  （完整交互在浏览器里——PPTX 内不存在 HTML 运行时，交互不伪嵌入）。

嵌入是构建后的独立步骤（build → capture → embed → 渲染验收）：
默认写 --output 新文件，不覆盖构建产物，便于对比与返工。

用法：
  python3 embed_animation.py --pptx deck.pptx --slide 5 \
      --gif out/demo.anim.gif --html out/demo.html \
      --box "0.6,1.4,8.2,4.6" --output deck.anim.pptx

- --box 单位英寸 "left,top,width,height"；缺省按幻灯片尺寸居中取 16:9 安全区，
  高度超过版心 62% 时自动回缩
- 图片在盒内等比缩放（fit-contain）并居中，不裁切不拉伸

依赖：python-pptx。缺包时给出与 compiler/build_pptx.py --ensure-deps 一致的指引。
工程约定：子进程一律参数数组（禁 shell 拼接）。
"""

from __future__ import annotations

import argparse
import os
import sys

EMU_PER_INCH = 914400


def fail(message: str, code: int = 1) -> None:
    print(f"[super-ppts] ERROR: {message}")
    raise SystemExit(code)


def ensure_pptx_module() -> None:
    try:
        import pptx  # noqa: F401
        return
    except ImportError:
        pass
    venv_python = os.path.expanduser(
        os.path.join("~", ".dsh", "venvs", "dsh-super-ppts", "bin", "python"))
    hints = [f"python3 <插件包根>/compiler/build_pptx.py --ensure-deps"]
    if os.path.isfile(venv_python):
        hints.append(f"{venv_python} {os.path.abspath(__file__)} …（专属 venv 内已有 python-pptx）")
    fail("python-pptx 不可用。任选其一：\n  " + "\n  ".join(hints))


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="动画素材嵌入器（GIF/快照 + 超链接）")
    parser.add_argument("--pptx", required=True, help="目标 PPTX（构建产物）")
    parser.add_argument("--slide", type=int, required=True, help="目标页序号（1-based）")
    parser.add_argument("--gif", help="动图 GIF（html_capture.py 的 .anim.gif；优先使用）")
    parser.add_argument("--snapshot", help="静态快照 PNG（无 --gif 时必填）")
    parser.add_argument("--html", help="随附交互 HTML 原件路径（作为图片超链接目标）")
    parser.add_argument("--box", metavar="L,T,W,H",
                        help="嵌入区域（英寸 left,top,width,height）；缺省居中 16:9 安全区")
    parser.add_argument("--note", help="可选说明文字（图下方小字，如「▶ 完整交互见随附 HTML」）")
    parser.add_argument("--link-absolute", action="store_true",
                        help="超链接用绝对路径（默认相对 PPTX 所在目录）")
    parser.add_argument("--output", help="输出 PPTX（缺省 <stem>.anim.pptx，与源同目录）")
    return parser.parse_args()


def default_box(slide_w_in: float, slide_h_in: float) -> tuple[float, float, float, float]:
    """居中 16:9 安全区：宽 80% 版面，高度触顶 62% 版面时回缩。"""
    width = slide_w_in * 0.80
    height = width * 9.0 / 16.0
    max_height = slide_h_in * 0.62
    if height > max_height:
        height = max_height
        width = height * 16.0 / 9.0
    left = (slide_w_in - width) / 2.0
    top = (slide_h_in - height) * 0.58
    return left, top, width, height


def parse_box(text: str | None, slide_w_in: float, slide_h_in: float) -> tuple[float, float, float, float]:
    if not text:
        return default_box(slide_w_in, slide_h_in)
    try:
        left, top, width, height = (float(part.strip()) for part in text.split(","))
    except ValueError:
        fail(f"--box 需要四个数字 L,T,W,H（英寸；收到：{text}）")
    if width <= 0 or height <= 0:
        fail(f"--box 宽高必须为正（收到：{text}）")
    return left, top, width, height


def resolve_link(html_path: str, pptx_path: str, absolute: bool) -> str:
    html_abs = os.path.abspath(html_path)
    if absolute:
        return html_abs
    rel = os.path.relpath(html_abs, start=os.path.dirname(os.path.abspath(pptx_path)))
    return rel.replace(os.sep, "/")  # PowerPoint 对正斜杠相对路径跨平台兼容


def main() -> None:
    args = parse_args()
    if not os.path.isfile(args.pptx):
        fail(f"PPTX 不存在：{args.pptx}")
    if args.slide < 1:
        fail(f"--slide 为 1-based 页序号（收到：{args.slide}）")
    media_path = args.gif or args.snapshot
    if media_path is None:
        fail("需要 --gif（优先）或 --snapshot（兜底）至少其一")
    for label, path in (("--gif", args.gif), ("--snapshot", args.snapshot)):
        if path and not os.path.isfile(path):
            fail(f"{label} 不存在：{path}")
    if args.html and not os.path.isfile(args.html):
        fail(f"--html 不存在：{args.html}")

    ensure_pptx_module()
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation(args.pptx)
    if args.slide > len(prs.slides):
        fail(f"--slide {args.slide} 超出页数（共 {len(prs.slides)} 页）")
    slide = prs.slides[args.slide - 1]

    slide_w_in = prs.slide_width / EMU_PER_INCH
    slide_h_in = prs.slide_height / EMU_PER_INCH
    left_in, top_in, box_w_in, box_h_in = parse_box(args.box, slide_w_in, slide_h_in)

    # 先按原生尺寸落图，读取纵横比后在盒内 fit-contain 并居中（不裁切不拉伸）
    picture = slide.shapes.add_picture(media_path, Inches(left_in), Inches(top_in))
    native_w = picture.width or 1
    native_h = picture.height or 1
    scale = min(Inches(box_w_in) / native_w, Inches(box_h_in) / native_h)
    picture.width = int(native_w * scale)
    picture.height = int(native_h * scale)
    picture.left = Inches(left_in) + (Inches(box_w_in) - picture.width) // 2
    picture.top = Inches(top_in) + (Inches(box_h_in) - picture.height) // 2

    link_target = None
    if args.html:
        link_target = resolve_link(args.html, args.pptx, args.link_absolute)
        picture.click_action.hyperlink.address = link_target

    if args.note:
        note_top = Inches(top_in + box_h_in)
        textbox = slide.shapes.add_textbox(Inches(left_in), note_top, Inches(box_w_in), Inches(0.3))
        frame = textbox.text_frame
        frame.text = args.note
        paragraph = frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(10)

    output = args.output or os.path.join(
        os.path.dirname(os.path.abspath(args.pptx)),
        os.path.splitext(os.path.basename(args.pptx))[0] + ".anim.pptx")
    prs.save(output)

    kind = "GIF 动图" if args.gif else "静态快照"
    print(f"[super-ppts] 嵌入完成：{kind} → 第 {args.slide} 页 "
          f"({picture.width / EMU_PER_INCH:.2f}x{picture.height / EMU_PER_INCH:.2f} in)")
    print(f"  素材: {media_path}")
    if link_target:
        print(f"  超链接: {link_target}{'（绝对）' if args.link_absolute else '（相对）'}")
    print(f"  输出: {output}")
    print("EMBED_OK")


if __name__ == "__main__":
    main()
